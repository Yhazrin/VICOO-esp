"""One-off: generate a single-slide tech summary PPTX. Run: python tools/generate_tech_onepager_pptx.py"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "tech-onepager.pptx"

prs = Presentation()
# 16:9 default; use blank layout (often index 6)
blank = 6
if len(prs.slide_layouts) > blank:
    slide_layout = prs.slide_layouts[blank]
else:
    slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)

# Title
left, top, width, height = Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.9)
box = slide.shapes.add_textbox(left, top, width, height)
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "Uniqlo × VICOO 公益·公益行动 — 技术架构（一页）"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
p.alignment = PP_ALIGN.LEFT

# Body
left, top, width, height = Inches(0.5), Inches(1.45), Inches(12.3), Inches(5.8)
body = slide.shapes.add_textbox(left, top, width, height)
bt = body.text_frame
bt.word_wrap = True

lines = [
    "整体：多端（Web 官网 / 管理台 / 微信小程序 / Android）共用一套 REST API。",
    "",
    "后端：FastAPI + 异步 SQLAlchemy，MySQL 持久化，Redis 缓存/会话，RabbitMQ 异步；JWT 等安全机制；业务含用户/活动/作品/电商/捐赠/供应链/管理/内容等。",
    "",
    "Web & 管理台：React 18 + Vite + TypeScript；TanStack Query + Zustand + Axios + i18n。",
    "",
    "部署：Docker Compose（MySQL + Redis + RabbitMQ + 后端 + Nginx 托管前端）· 敏感信息走 Secrets。",
    "",
    "一句话：React/Vite 多端 + FastAPI 单体 API，MySQL/Redis 支撑，容器化一体部署。",
]

for i, line in enumerate(lines):
    if i == 0:
        para = bt.paragraphs[0]
    else:
        para = bt.add_paragraph()
    para.text = line
    para.font.size = Pt(16) if line else Pt(8)
    para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    para.space_after = Pt(4)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"Wrote: {OUT}")
